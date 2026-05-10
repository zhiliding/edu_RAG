from typing import Iterator
from edu_ocr import get_ocr
from langchain_core.documents import Document
from langchain_core.document_loaders import BaseLoader
from pptx import Presentation
from PIL import Image
import numpy as np
from io import BytesIO
from tqdm import tqdm


class OCRPPTLoader(BaseLoader):
    """An example document loader that reads a file line by line."""

    def __init__(self, filepath: str) -> None:
        """Initialize the loader with a file path.

        Args:
            filepath: The path to the ppt to load.
        """
        self.filepath = filepath

    def lazy_load(self) -> Iterator[Document]:
        # <-- Does not take any arguments
        """A lazy loader that reads a file line by line.

        When you're implementing lazy load methods, you should use a generator
        to yield documents one by one.
        """

        line = self.ppt2text(self.filepath)
        yield Document(page_content=line, metadata={"source": self.filepath})

    def ppt2text(self, filepath):
        # 打开指定路径的 PowerPoint 文件
        prs = Presentation(filepath)
        # print(f'prs-->{prs}')
        # 获取 OCR 功能的实例
        ocr = get_ocr()
        # 初始化一个空字符串，用于存储提取的文本内容
        resp = ""

        def extract_text(shape):
            # nonlocal指明resp非全局非局部，而是外部嵌套函数中的变量，
            # 允许内部函数访问和修改外部函数中定义的变量resp
            nonlocal resp

            # 检查形状是否有文本框
            if shape.has_text_frame:
                # 将文本框中的文本添加到resp中，并去掉前后空格
                resp += shape.text.strip() + "\n"

            # 检查形状是否为表格
            if shape.has_table:
                # 遍历表格的每一行
                for row in shape.table.rows:
                    # 遍历每一行中的每个单元格
                    for cell in row.cells:
                        # 遍历单元格中的每个段落
                        for paragraph in cell.text_frame.paragraphs:
                            # 将单元格中的文本添加到resp中，并去掉前后空格
                            resp += paragraph.text.strip() + "\n"

            # 检查形状是否为图片（shape_type == 13）
            if shape.shape_type == 13:  # 13 表示图片
                # 使用 BytesIO 打开图片数据并转换为图像对象
                image = Image.open(BytesIO(shape.image.blob))
                # 使用 OCR 处理图像并获取结果
                result, _ = ocr(np.array(image))
                if result:  # 如果 OCR 有结果
                    # 提取 OCR 结果中的文本行
                    ocr_result = [line[1] for line in result]
                    # 将 OCR 提取的文本添加到resp中，以换行分隔
                    resp += "\n".join(ocr_result)

            # 检查形状是否为组合形状（shape_type == 6）
            elif shape.shape_type == 6:  # 6 表示组合
                # 遍历组合形状中的每个子形状，递归调用extract_text函数
                for child_shape in shape.shapes:
                    extract_text(child_shape)

        # 创建一个进度条，用于显示幻灯片处理进度，初始总数为幻灯片数量
        b_unit = tqdm(total=len(prs.slides), desc="OCRPPTLoader slide index: 1")

        # 遍历所有幻灯片
        for slide_number, slide in enumerate(prs.slides, start=1):
            # 更新进度条描述，显示当前处理的幻灯片索引
            b_unit.set_description("OCRPPTLoader slide index: {}".format(slide_number))
            b_unit.refresh()  # 刷新进度条显示

            # 按照从上到下、从左到右的顺序对形状进行排序遍历
            sorted_shapes = sorted(slide.shapes, key=lambda x: (x.top, x.left))

            for shape in sorted_shapes:
                extract_text(shape)  # 调用extract_text函数提取当前形状的文本内容

            b_unit.update(1)  # 更新进度条，表示处理了一张幻灯片

        return resp  # 返回提取到的所有文本内容


if __name__ == '__main__':
    img_loader = OCRPPTLoader(filepath='/Users/ligang/Desktop/AI29期课堂资料/Codes/integrated_qa_system/rag_qa/samples/ocr_01.pptx')
    doc = img_loader.load()
    print(doc)